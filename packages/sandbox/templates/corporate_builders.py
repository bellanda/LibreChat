"""Reusable enterprise-clean builders for dynamic artifact generation.

Use these helpers when the assistant needs to generate business files
with consistent HPE Corporate Clean visual style across all formats.

Supports: PDF (Jinja2), PPTX (6 slide types), XLSX (multi-sheet with charts), DOCX (programmatic).
"""

from __future__ import annotations

from pathlib import Path
from typing import Any

import polars as pl
from jinja2 import Environment, FileSystemLoader
from xhtml2pdf import pisa
from pptx import Presentation
from pptx.util import Inches, Pt as PptPt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor as PptxRGBColor
from openpyxl import Workbook
from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
from openpyxl.chart import BarChart, Reference
from docx import Document
from docx.shared import Pt, RGBColor, Inches as DocxInches
from docx.enum.text import WD_PARAGRAPH_ALIGNMENT

from style import get_theme

TEMPLATES_DIR = Path(__file__).parent / "layouts"


def _resolve_css_vars(html: str, theme_cfg) -> str:
    """Replace CSS custom properties (var(--x)) with their real hex values.

    xhtml2pdf does not support var(). This substitution runs on the rendered
    HTML so the PDF renderer sees only concrete color values.
    """
    import re
    replacements = {
        "var(--primary)": f"#{theme_cfg.primary}",
        "var(--secondary)": f"#{theme_cfg.secondary}",
        "var(--accent)": f"#{theme_cfg.accent}",
        "var(--accent-light)": f"#{theme_cfg.accent_light}",
        "var(--success)": f"#{theme_cfg.success}",
        "var(--text)": f"#{theme_cfg.text}",
        "var(--muted)": f"#{theme_cfg.muted}",
        "var(--bg-light)": f"#{theme_cfg.surface}",
    }
    for var, value in replacements.items():
        html = html.replace(var, value)
    return html


def _pptx_rgb(hex_color: str) -> PptxRGBColor:
    """Convert hex color to pptx RGB color."""
    hex_clean = hex_color.lstrip("#")
    return PptxRGBColor.from_string(hex_clean)


def _style_pptx_slide_with_title(slide, theme_cfg, title: str, subtitle: str = ""):
    """Add HPE-styled title bar to a slide."""
    bg_shape = slide.shapes.add_shape(
        1,  # RECTANGLE
        Inches(0),
        Inches(0),
        Inches(10),
        Inches(1.2),
    )
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = _pptx_rgb(theme_cfg.primary)
    bg_shape.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.7))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_p = title_frame.paragraphs[0]
    title_p.font.name = theme_cfg.font_main
    title_p.font.size = PptPt(32)
    title_p.font.bold = True
    title_p.font.color.rgb = _pptx_rgb("FFFFFF")

    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(9), Inches(0.4))
        sub_frame = sub_box.text_frame
        sub_frame.text = subtitle
        sub_p = sub_frame.paragraphs[0]
        sub_p.font.name = theme_cfg.font_main
        sub_p.font.size = PptPt(14)
        sub_p.font.color.rgb = _pptx_rgb(theme_cfg.accent)


# ── PPTX auto-slide generation ────────────────────────────────────────────────

def _auto_slides_from_document_type(document_type: str, kw: dict) -> list[dict]:
    """Convert document_type + kwargs into a list of slide specs."""
    title = kw.get("title", "")
    date = kw.get("date", "")
    subtitle = kw.get("subtitle", "")
    slides: list[dict] = [{"type": "cover", "title": title, "subtitle": subtitle, "date": date}]

    if document_type in (None, "relatorio"):
        if kw.get("kpis"):
            slides.append({"type": "kpi_grid", "title": "Indicadores", "kpis": kw["kpis"]})
        for section in kw.get("sections", []):
            tbl = section.get("table")
            if tbl:
                slides.append({"type": "table", "title": section.get("title", ""),
                                "headers": tbl.get("headers", []), "rows": tbl.get("rows", [])})
            elif section.get("bullets"):
                slides.append({"type": "bullets", "title": section.get("title", ""),
                                "bullets": section["bullets"]})
            elif section.get("content"):
                slides.append({"type": "bullets", "title": section.get("title", ""),
                                "bullets": [section["content"]]})

    elif document_type == "tutorial":
        for step in kw.get("steps", []):
            bullets = [step.get("content", "")]
            if step.get("code"):
                bullets.append(f"Código: {step['code'][:120]}")
            slides.append({"type": "bullets", "title": step.get("title", ""), "bullets": bullets})
        if kw.get("tips"):
            slides.append({"type": "bullets", "title": "Dicas", "bullets": kw["tips"]})
        if kw.get("warnings"):
            slides.append({"type": "bullets", "title": "Atenção", "bullets": kw["warnings"]})

    elif document_type == "ata":
        attendees = kw.get("attendees", [])
        if attendees:
            slides.append({"type": "bullets", "title": "Participantes", "bullets": attendees})
        decisions = kw.get("decisions", [])
        if decisions:
            slides.append({"type": "bullets", "title": "Decisões", "bullets": decisions})
        action_items = kw.get("action_items", [])
        if action_items:
            rows = [[ai.get("task", ""), ai.get("owner", "—"), ai.get("due", "—")] for ai in action_items]
            slides.append({"type": "table", "title": "Action Items",
                           "headers": ["Tarefa", "Responsável", "Prazo"], "rows": rows})

    elif document_type == "brainstorm":
        ideas = kw.get("ideas", [])
        if ideas:
            rows = [[i.get("title", ""), i.get("category", "—"), i.get("priority", "—"),
                     i.get("description", "")] for i in ideas]
            slides.append({"type": "table", "title": "Ideias",
                           "headers": ["Ideia", "Categoria", "Prioridade", "Descrição"], "rows": rows})
        if kw.get("next_steps"):
            slides.append({"type": "bullets", "title": "Próximos Passos", "bullets": kw["next_steps"]})

    elif document_type == "proposta":
        scope = kw.get("scope", "")
        deliverables = kw.get("deliverables", [])
        bullets = ([scope] if scope else []) + (deliverables if isinstance(deliverables, list) else [deliverables])
        if bullets:
            slides.append({"type": "bullets", "title": "Escopo e Entregas", "bullets": bullets})
        pricing = kw.get("pricing", [])
        if pricing and isinstance(pricing, list) and isinstance(pricing[0], dict):
            rows = [[p.get("item", ""), p.get("description", ""), p.get("value", "")] for p in pricing]
            slides.append({"type": "table", "title": "Investimento",
                           "headers": ["Item", "Descrição", "Valor"], "rows": rows})

    elif document_type == "spec_tecnica":
        requirements = kw.get("requirements", [])
        if requirements:
            if isinstance(requirements[0], dict):
                rows = [[str(r.get("id", i + 1)), r.get("description", ""), r.get("priority", "—")]
                        for i, r in enumerate(requirements)]
                slides.append({"type": "table", "title": "Requisitos",
                               "headers": ["ID", "Requisito", "Prioridade"], "rows": rows})
            else:
                slides.append({"type": "bullets", "title": "Requisitos", "bullets": requirements})
        decisions = kw.get("decisions", [])
        if decisions and isinstance(decisions[0], dict):
            rows = [[d.get("decision", ""), d.get("rationale", ""), d.get("status", "")] for d in decisions]
            slides.append({"type": "table", "title": "Decisões de Design",
                           "headers": ["Decisão", "Racional", "Status"], "rows": rows})

    elif document_type == "plano_projeto":
        milestones = kw.get("milestones", [])
        if milestones:
            slides.append({"type": "timeline", "title": "Marcos do Projeto",
                           "milestones": [{"label": m.get("label", ""), "text": m.get("date", "")}
                                          for m in milestones]})
            rows = [[m.get("label", ""), m.get("owner", "—"), m.get("date", "—"), m.get("status", "—")]
                    for m in milestones]
            slides.append({"type": "table", "title": "Detalhes",
                           "headers": ["Marco", "Responsável", "Data", "Status"], "rows": rows})
        risks = kw.get("risks", [])
        if risks:
            rows = [[r.get("risk", ""), r.get("probability", "—"), r.get("impact", "—"),
                     r.get("mitigation", "—")] for r in risks]
            slides.append({"type": "table", "title": "Riscos",
                           "headers": ["Risco", "Probabilidade", "Impacto", "Mitigação"], "rows": rows})

    elif document_type == "onboarding":
        steps = kw.get("steps", [])
        if steps:
            bullets = [f"{s.get('title', '')}: {s.get('content', '')}" for s in steps]
            slides.append({"type": "bullets", "title": "Passos de Integração", "bullets": bullets})
        contacts = kw.get("contacts", [])
        if contacts:
            rows = [[c.get("name", ""), c.get("role", ""), c.get("contact", "")] for c in contacts]
            slides.append({"type": "table", "title": "Contatos",
                           "headers": ["Nome", "Papel", "Contato"], "rows": rows})
        checklist = kw.get("checklist", [])
        if checklist:
            slides.append({"type": "bullets", "title": "Checklist", "bullets": checklist})

    return slides


# ── DOCX data normalization ───────────────────────────────────────────────────

def _normalize_docx_data(document_type: str | None, kwargs: dict) -> dict:
    """Convert document_type + raw kwargs into canonical DOCX data dict."""
    data: dict = dict(kwargs)
    data.setdefault("document_title", data.pop("title", "Documento"))

    if document_type in (None, "relatorio"):
        return data

    sections: list[dict] = []

    if document_type == "tutorial":
        for step in data.get("steps", []):
            sec: dict = {"title": step.get("title", ""), "content": step.get("content", "")}
            if step.get("code"):
                sec["bullets"] = [f"Código: {step['code']}"]
            sections.append(sec)
        if data.get("tips"):
            sections.append({"title": "Dicas", "bullets": data["tips"]})
        if data.get("warnings"):
            sections.append({"title": "Atenção", "bullets": data["warnings"]})

    elif document_type == "ata":
        if data.get("attendees"):
            sections.append({"title": "Participantes", "bullets": data["attendees"]})
        if data.get("decisions"):
            sections.append({"title": "Decisões", "bullets": data["decisions"]})
        action_items = data.get("action_items", [])
        if action_items:
            rows = [[ai.get("task", ""), ai.get("owner", "—"), ai.get("due", "—")] for ai in action_items]
            sections.append({"title": "Action Items",
                              "table": {"headers": ["Tarefa", "Responsável", "Prazo"], "rows": rows}})

    elif document_type == "brainstorm":
        if data.get("context"):
            sections.append({"title": "Contexto", "content": data["context"]})
        ideas = data.get("ideas", [])
        if ideas:
            rows = [[i.get("title", ""), i.get("category", "—"), i.get("priority", "—"),
                     i.get("description", "")] for i in ideas]
            sections.append({"title": "Ideias",
                              "table": {"headers": ["Ideia", "Categoria", "Prioridade", "Descrição"],
                                        "rows": rows}})
        if data.get("next_steps"):
            sections.append({"title": "Próximos Passos", "bullets": data["next_steps"]})

    elif document_type == "proposta":
        if data.get("scope"):
            sections.append({"title": "Escopo", "content": data["scope"]})
        deliverables = data.get("deliverables", [])
        if deliverables:
            bullets = deliverables if isinstance(deliverables, list) else [deliverables]
            sections.append({"title": "Entregas", "bullets": bullets})
        pricing = data.get("pricing", [])
        if pricing and isinstance(pricing, list) and isinstance(pricing[0], dict):
            rows = [[p.get("item", ""), p.get("description", ""), p.get("value", "")] for p in pricing]
            sections.append({"title": "Investimento",
                              "table": {"headers": ["Item", "Descrição", "Valor"], "rows": rows}})

    elif document_type == "spec_tecnica":
        requirements = data.get("requirements", [])
        if requirements:
            if isinstance(requirements[0], dict):
                rows = [[str(r.get("id", i + 1)), r.get("description", ""), r.get("priority", "—")]
                        for i, r in enumerate(requirements)]
                sections.append({"title": "Requisitos",
                                  "table": {"headers": ["ID", "Requisito", "Prioridade"], "rows": rows}})
            else:
                sections.append({"title": "Requisitos", "bullets": requirements})
        decisions = data.get("decisions", [])
        if decisions and isinstance(decisions[0], dict):
            rows = [[d.get("decision", ""), d.get("rationale", ""), d.get("status", "")] for d in decisions]
            sections.append({"title": "Decisões de Design",
                              "table": {"headers": ["Decisão", "Racional", "Status"], "rows": rows}})

    elif document_type == "plano_projeto":
        milestones = data.get("milestones", [])
        if milestones:
            rows = [[m.get("label", ""), m.get("owner", "—"), m.get("date", "—"), m.get("status", "—")]
                    for m in milestones]
            sections.append({"title": "Marcos",
                              "table": {"headers": ["Marco", "Responsável", "Data", "Status"], "rows": rows}})
        risks = data.get("risks", [])
        if risks:
            rows = [[r.get("risk", ""), r.get("probability", "—"), r.get("impact", "—"),
                     r.get("mitigation", "—")] for r in risks]
            sections.append({"title": "Riscos",
                              "table": {"headers": ["Risco", "Probabilidade", "Impacto", "Mitigação"],
                                        "rows": rows}})

    elif document_type == "onboarding":
        for step in data.get("steps", []):
            sections.append({"title": step.get("title", ""), "content": step.get("content", "")})
        contacts = data.get("contacts", [])
        if contacts:
            rows = [[c.get("name", ""), c.get("role", ""), c.get("contact", "")] for c in contacts]
            sections.append({"title": "Contatos",
                              "table": {"headers": ["Nome", "Papel", "Contato"], "rows": rows}})
        if data.get("checklist"):
            sections.append({"title": "Checklist", "bullets": data["checklist"]})

    data["sections"] = sections
    return data


# ── Public builders ───────────────────────────────────────────────────────────

def create_pptx_document(
    output_path: str | Path,
    *,
    theme: str = "executivo",
    slides: list[dict[str, Any]] | None = None,
    document_type: str | None = None,
    **kwargs: Any,
) -> Path:
    """Create a PPTX presentation with 6 slide types.

    Slide types: cover, kpi_grid, bullets, chart, table, timeline.

    Can be called two ways:
    - Explicit: pass ``slides=[...]``
    - Auto: pass ``document_type="tutorial"`` + document fields as kwargs
    """
    theme_cfg = get_theme(theme)
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    if slides is None and document_type is not None:
        slides = _auto_slides_from_document_type(document_type, kwargs)

    for slide_spec in slides or []:
        slide_type = slide_spec.get("type", "bullets")

        if slide_type == "cover":
            _add_cover_slide(prs, theme_cfg, slide_spec)
        elif slide_type == "kpi_grid":
            _add_kpi_grid_slide(prs, theme_cfg, slide_spec)
        elif slide_type == "bullets":
            _add_bullets_slide(prs, theme_cfg, slide_spec)
        elif slide_type == "chart":
            _add_chart_slide(prs, theme_cfg, slide_spec)
        elif slide_type == "table":
            _add_table_slide(prs, theme_cfg, slide_spec)
        elif slide_type == "timeline":
            _add_timeline_slide(prs, theme_cfg, slide_spec)

    out = Path(output_path)
    prs.save(str(out))
    return out


def _add_cover_slide(prs, theme_cfg, spec: dict):
    """Create cover slide with HPE blue background."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = _pptx_rgb(theme_cfg.primary)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.text = spec.get("title", "")
    title_frame.word_wrap = True
    for para in title_frame.paragraphs:
        para.font.size = PptPt(54)
        para.font.bold = True
        para.font.color.rgb = _pptx_rgb("FFFFFF")
        para.font.name = theme_cfg.font_main

    # Subtitle
    if spec.get("subtitle"):
        sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.1), Inches(9), Inches(1))
        sub_frame = sub_box.text_frame
        sub_frame.text = spec.get("subtitle")
        for para in sub_frame.paragraphs:
            para.font.size = PptPt(28)
            para.font.color.rgb = _pptx_rgb(theme_cfg.accent)
            para.font.name = theme_cfg.font_main

    # Date
    if spec.get("date"):
        date_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.5), Inches(9), Inches(0.5))
        date_frame = date_box.text_frame
        date_frame.text = spec.get("date")
        for para in date_frame.paragraphs:
            para.font.size = PptPt(14)
            para.font.color.rgb = _pptx_rgb(theme_cfg.muted)
            para.font.name = theme_cfg.font_main


def _add_kpi_grid_slide(prs, theme_cfg, spec: dict):
    """Create KPI grid slide (2x2 layout)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    _style_pptx_slide_with_title(slide, theme_cfg, spec.get("title", "KPIs"))

    kpis = spec.get("kpis", [])
    positions = [
        (0.5, 1.5), (5.25, 1.5),
        (0.5, 4.5), (5.25, 4.5),
    ]

    for idx, kpi in enumerate(kpis[:4]):
        x, y = positions[idx]
        box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(4.25), Inches(2.5))
        tf = box.text_frame

        # Value
        p_val = tf.paragraphs[0]
        p_val.text = kpi.get("value", "N/A")
        p_val.font.size = PptPt(36)
        p_val.font.bold = True
        p_val.font.color.rgb = _pptx_rgb(theme_cfg.primary)
        p_val.alignment = PP_ALIGN.CENTER

        # Label
        p_label = tf.add_paragraph()
        p_label.text = kpi.get("label", "")
        p_label.font.size = PptPt(12)
        p_label.font.color.rgb = _pptx_rgb(theme_cfg.muted)
        p_label.alignment = PP_ALIGN.CENTER


def _add_bullets_slide(prs, theme_cfg, spec: dict):
    """Create bullets slide with section header bar."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _style_pptx_slide_with_title(slide, theme_cfg, spec.get("title", ""))

    bullets = spec.get("bullets", [])
    content_box = slide.shapes.add_textbox(Inches(1), Inches(1.8), Inches(8), Inches(5))
    tf = content_box.text_frame
    tf.word_wrap = True

    for idx, bullet in enumerate(bullets):
        if idx == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {bullet}"
        p.font.size = PptPt(18)
        p.font.color.rgb = _pptx_rgb(theme_cfg.text)
        p.font.name = theme_cfg.font_main
        p.level = 0


def _add_chart_slide(prs, theme_cfg, spec: dict):
    """Create chart slide with embedded PNG image."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _style_pptx_slide_with_title(slide, theme_cfg, spec.get("title", ""))

    if spec.get("chart_path"):
        try:
            slide.shapes.add_picture(
                spec.get("chart_path"),
                Inches(1),
                Inches(1.8),
                width=Inches(8),
            )
        except Exception as e:
            print(f"Error adding chart image: {e}")


def _add_table_slide(prs, theme_cfg, spec: dict):
    """Create table slide with HPE-styled table."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _style_pptx_slide_with_title(slide, theme_cfg, spec.get("title", ""))

    headers = spec.get("headers", [])
    rows = spec.get("rows", [])

    if not headers:
        return

    table_shape = slide.shapes.add_table(
        len(rows) + 1, len(headers),
        Inches(0.5), Inches(1.8), Inches(9), Inches(4.5),
    )
    table = table_shape.table

    # Header row — colored background + white text
    for col_idx, header in enumerate(headers):
        cell = table.cell(0, col_idx)
        cell.fill.solid()
        cell.fill.fore_color.rgb = _pptx_rgb(theme_cfg.primary)
        tf = cell.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = str(header)
        run.font.bold = True
        run.font.color.rgb = _pptx_rgb("FFFFFF")
        run.font.size = PptPt(11)
        run.font.name = theme_cfg.font_main

    # Body rows
    for row_idx, row_data in enumerate(rows, 1):
        for col_idx, cell_value in enumerate(row_data):
            cell = table.cell(row_idx, col_idx)
            tf = cell.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = str(cell_value)
            run.font.size = PptPt(10)
            run.font.name = theme_cfg.font_main


def _add_timeline_slide(prs, theme_cfg, spec: dict):
    """Create timeline slide with milestone markers."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _style_pptx_slide_with_title(slide, theme_cfg, spec.get("title", ""))

    milestones = spec.get("milestones", [])
    if not milestones:
        return

    # Draw timeline line
    line = slide.shapes.add_connector(1, Inches(1), Inches(3.5), Inches(9), Inches(3.5))
    line.line.color.rgb = _pptx_rgb(theme_cfg.primary)
    line.line.width = PptPt(2)

    # Add milestone markers
    step = 8 / len(milestones)
    for idx, milestone in enumerate(milestones):
        x = 1 + (step * idx)

        # Dot
        dot = slide.shapes.add_shape(1, Inches(x - 0.15), Inches(3.35), Inches(0.3), Inches(0.3))
        dot.fill.solid()
        dot.fill.fore_color.rgb = _pptx_rgb(theme_cfg.primary)
        dot.line.color.rgb = _pptx_rgb("FFFFFF")
        dot.line.width = PptPt(2)

        # Label
        label_box = slide.shapes.add_textbox(Inches(x - 0.5), Inches(4), Inches(1), Inches(0.5))
        label_frame = label_box.text_frame
        label_frame.text = milestone.get("label", "")
        for para in label_frame.paragraphs:
            para.font.size = PptPt(10)
            para.font.bold = True
            para.alignment = PP_ALIGN.CENTER
            para.font.color.rgb = _pptx_rgb(theme_cfg.text)

        # Description
        if milestone.get("text"):
            desc_box = slide.shapes.add_textbox(Inches(x - 0.7), Inches(4.7), Inches(1.4), Inches(1.5))
            desc_frame = desc_box.text_frame
            desc_frame.word_wrap = True
            desc_frame.text = milestone.get("text")
            for para in desc_frame.paragraphs:
                para.font.size = PptPt(9)
                para.alignment = PP_ALIGN.CENTER
                para.font.color.rgb = _pptx_rgb(theme_cfg.muted)


def create_xlsx_document(
    output_path: str | Path,
    *,
    theme: str = "executivo",
    sheets: list[dict[str, Any]] | None = None,
) -> Path:
    """Create multi-sheet XLSX with charts and HPE styling.

    Each sheet can be:
    - type: 'dashboard' (KPIs + data table + chart)
    - type: 'data' (simple data table)
    - type: 'chart' (data + embedded bar chart)
    """
    theme_cfg = get_theme(theme)
    wb = Workbook()
    wb.remove(wb.active)  # Remove default sheet

    for sheet_spec in sheets or []:
        sheet_name = sheet_spec.get("name", "Sheet")[:31]  # Max 31 chars in Excel
        sheet_type = sheet_spec.get("type", "data")

        if sheet_type == "dashboard":
            _add_dashboard_sheet(wb, theme_cfg, sheet_name, sheet_spec)
        elif sheet_type == "data":
            _add_data_sheet(wb, theme_cfg, sheet_name, sheet_spec)
        elif sheet_type == "chart":
            _add_chart_sheet(wb, theme_cfg, sheet_name, sheet_spec)

    out = Path(output_path)
    wb.save(str(out))
    return out


def _add_dashboard_sheet(wb, theme_cfg, sheet_name: str, spec: dict):
    """Create dashboard sheet with KPIs, data, and optional chart."""
    ws = wb.create_sheet(sheet_name)
    row = 1

    # Title
    if spec.get("title"):
        ws[f"A{row}"] = spec.get("title")
        ws[f"A{row}"].font = Font(name=theme_cfg.font_main, size=14, bold=True, color=theme_cfg.primary)
        row += 2

    # KPIs
    if spec.get("kpis"):
        for kpi in spec.get("kpis", []):
            ws[f"A{row}"] = kpi.get("label", "")
            ws[f"B{row}"] = kpi.get("value", "")
            if kpi.get("delta"):
                ws[f"C{row}"] = kpi.get("delta")
            row += 1
        row += 1

    # Data table
    headers = spec.get("headers", [])
    rows_data = spec.get("rows", [])

    if headers:
        for col_idx, header in enumerate(headers, 1):
            cell = ws.cell(row, col_idx)
            cell.value = header
            cell.font = Font(name=theme_cfg.font_main, bold=True, color="FFFFFF", size=10)
            cell.fill = PatternFill(start_color=theme_cfg.primary, end_color=theme_cfg.primary, fill_type="solid")
            cell.alignment = Alignment(horizontal="center", vertical="center")

        # Data rows
        for row_idx, row_data in enumerate(rows_data, row + 1):
            for col_idx, cell_value in enumerate(row_data, 1):
                cell = ws.cell(row_idx, col_idx)
                cell.value = cell_value
                cell.font = Font(name=theme_cfg.font_main, size=9)
                # Alternating row colors
                if row_idx % 2 == 0:
                    cell.fill = PatternFill(start_color="F9FAFC", end_color="F9FAFC", fill_type="solid")

        # Add chart if requested
        if spec.get("add_chart") and len(rows_data) > 0:
            chart = BarChart()
            chart.title = spec.get("chart_title", "Chart")
            data = Reference(ws, min_col=2, min_row=row, max_row=row + len(rows_data))
            chart.add_data(data)
            ws.add_chart(chart, f"A{row + len(rows_data) + 3}")

    # Auto column width
    for col_idx, header in enumerate(headers, 1):
        ws.column_dimensions[chr(64 + col_idx)].width = 18


def _add_data_sheet(wb, theme_cfg, sheet_name: str, spec: dict):
    """Create simple data table sheet."""
    ws = wb.create_sheet(sheet_name)
    row = 1

    # Title
    if spec.get("title"):
        ws[f"A{row}"] = spec.get("title")
        ws[f"A{row}"].font = Font(name=theme_cfg.font_main, size=12, bold=True, color=theme_cfg.primary)
        row += 2

    # Headers
    headers = spec.get("headers", [])
    rows_data = spec.get("rows", [])

    for col_idx, header in enumerate(headers, 1):
        cell = ws.cell(row, col_idx)
        cell.value = header
        cell.font = Font(name=theme_cfg.font_main, bold=True, color="FFFFFF", size=10)
        cell.fill = PatternFill(start_color=theme_cfg.primary, end_color=theme_cfg.primary, fill_type="solid")
        cell.alignment = Alignment(horizontal="center", vertical="center")

    # Data rows
    for row_idx, row_data in enumerate(rows_data, row + 1):
        for col_idx, cell_value in enumerate(row_data, 1):
            cell = ws.cell(row_idx, col_idx)
            cell.value = cell_value
            cell.font = Font(name=theme_cfg.font_main, size=9)
            if row_idx % 2 == 0:
                cell.fill = PatternFill(start_color="F9FAFC", end_color="F9FAFC", fill_type="solid")

    # Auto column width
    for col_idx in range(1, len(headers) + 1):
        ws.column_dimensions[chr(64 + col_idx)].width = 18


def _add_chart_sheet(wb, theme_cfg, sheet_name: str, spec: dict):
    """Create sheet with data and embedded chart."""
    _add_dashboard_sheet(wb, theme_cfg, sheet_name, {**spec, "add_chart": True})


def create_pdf_document(
    output_path: str | Path,
    *,
    template: str | None = None,
    theme: str = "executivo",
    data: dict[str, Any] | None = None,
    document_type: str | None = None,
    **kwargs: Any,
) -> Path:
    """Create PDF using Jinja2 templates + xhtml2pdf.

    Can be called two ways:
    - Old style: ``data={"document_title": ..., "sections": [...]}``
    - New style: pass ``document_type`` and document fields as kwargs

    Available themes: executivo, operacional, tecnico
    """
    out = Path(output_path)

    # Build template context
    if data is not None:
        ctx = dict(data)
    else:
        ctx = dict(kwargs)

    # Inject document_type so the template can branch
    if document_type is not None:
        ctx.setdefault("document_type", document_type)

    # Layouts templates use `title`; old pdf_templates used `document_title` — normalise both
    if "title" not in ctx and "document_title" in ctx:
        ctx["title"] = ctx["document_title"]
    if "document_title" not in ctx and "title" in ctx:
        ctx["document_title"] = ctx["title"]
    ctx.setdefault("title", "Documento")
    ctx.setdefault("document_title", ctx["title"])

    theme_cfg = get_theme(theme)

    # Determine which template file to use
    tpl_name = f"{template or theme}.html"
    env = Environment(loader=FileSystemLoader(str(TEMPLATES_DIR)))
    try:
        jinja_template = env.get_template(tpl_name)
    except Exception as e:
        raise ValueError(f"Template '{tpl_name}' not found in {TEMPLATES_DIR}: {e}")

    html_content = jinja_template.render(**ctx)
    html_content = _resolve_css_vars(html_content, theme_cfg)

    with open(out, "w+b") as pdf_file:
        pisa.CreatePDF(html_content, pdf_file)

    return out


def create_docx_document(
    output_path: str | Path,
    *,
    theme: str = "executivo",
    data: dict[str, Any] | None = None,
    document_type: str | None = None,
    **kwargs: Any,
) -> Path:
    """Create DOCX programmatically with HPE styling.

    Can be called two ways:
    - Old style: ``data={"document_title": ..., "sections": [...]}``
    - New style: pass ``document_type`` and document fields as kwargs
    """
    theme_cfg = get_theme(theme)

    # Resolve data dict
    if data is not None:
        doc_data = dict(data)
    else:
        doc_data = _normalize_docx_data(document_type, kwargs)

    doc_data.setdefault("document_title", doc_data.pop("title", "Documento"))

    doc = Document()

    # Set default font
    style = doc.styles["Normal"]
    font = style.font
    font.name = theme_cfg.font_main
    font.size = Pt(11)
    font.color.rgb = RGBColor.from_string(theme_cfg.text)

    # Document title
    if doc_data.get("document_title"):
        title = doc.add_paragraph()
        title_run = title.add_run(doc_data.get("document_title"))
        title_run.font.size = Pt(24)
        title_run.font.bold = True
        title_run.font.color.rgb = RGBColor.from_string(theme_cfg.primary)
        title.alignment = WD_PARAGRAPH_ALIGNMENT.CENTER

    # Author / Date / Classification
    if doc_data.get("author") or doc_data.get("date") or doc_data.get("classification"):
        meta = doc.add_paragraph()
        if doc_data.get("author"):
            meta.add_run(f"Autor: {doc_data.get('author')} | ")
        if doc_data.get("date"):
            meta.add_run(f"Data: {doc_data.get('date')} | ")
        if doc_data.get("classification"):
            meta.add_run(f"Classificação: {doc_data.get('classification')}")
        for run in meta.runs:
            run.font.size = Pt(9)
            run.font.color.rgb = RGBColor.from_string(theme_cfg.muted)
        meta.paragraph_format.space_after = Pt(12)

    # Abstract
    if doc_data.get("abstract") or doc_data.get("executive_summary"):
        text = doc_data.get("abstract") or doc_data.get("executive_summary")
        abstract = doc.add_paragraph(text)
        abstract.paragraph_format.left_indent = DocxInches(0.25)
        for run in abstract.runs:
            run.font.italic = True
            run.font.color.rgb = RGBColor.from_string(theme_cfg.secondary)
        abstract.paragraph_format.space_after = Pt(12)

    # Sections
    for section in doc_data.get("sections", []):
        if section.get("title"):
            doc.add_heading(section.get("title"), level=1)

        if section.get("content"):
            doc.add_paragraph(section.get("content"))

        if section.get("bullets"):
            for bullet in section["bullets"]:
                p = doc.add_paragraph(style="List Bullet")
                p.add_run(str(bullet))

        if section.get("table"):
            table_spec = section.get("table")
            headers = table_spec.get("headers", [])
            rows = table_spec.get("rows", [])

            if headers:
                table = doc.add_table(rows=len(rows) + 1, cols=len(headers))
                table.style = "Light Grid Accent 1"

                header_cells = table.rows[0].cells
                for col_idx, header in enumerate(headers):
                    cell = header_cells[col_idx]
                    cell.text = str(header)
                    for paragraph in cell.paragraphs:
                        for run in paragraph.runs:
                            run.font.bold = True

                for row_idx, row_data in enumerate(rows, 1):
                    row_cells = table.rows[row_idx].cells
                    for col_idx, cell_value in enumerate(row_data):
                        row_cells[col_idx].text = str(cell_value)

    # Page header/footer
    section = doc.sections[0]
    header = section.header
    header_para = header.paragraphs[0]
    header_para.text = "HPE Intelligence Platform"

    footer = section.footer
    footer_para = footer.paragraphs[0]
    footer_para.text = f"HPE Confidential | {doc_data.get('document_title', 'Document')}"

    out = Path(output_path)
    doc.save(str(out))
    return out


def create_csv_document(output_path: str | Path, *, rows: list[dict[str, Any]]) -> Path:
    """Create simple CSV file from row data."""
    out = Path(output_path)
    pl.DataFrame(rows if rows else [{"coluna": "sem_dados"}]).write_csv(str(out))
    return out
